from flask import Flask, render_template, request, send_file
from bs4 import BeautifulSoup
import requests
import re
from fpdf import FPDF
from pptx import Presentation
from docx import Document
from pptx.util import Inches

app = Flask(__name__)

# Function to scrape information from different sources
def scrape_information(query, query_type):
    # Byju's URL
    byjus_url = f"https://www.byjus.com/search/{query}"
    # Other sources URLs
    other_sources = [
        f"https://scholar.google.com/scholar?q={query}",
        f"https://en.wikipedia.org/wiki/{query}",
        f"https://sw.wikipedia.org/wiki/{query}",
        f"https://www.britannica.com/wiki/{query}"
    ]
    content_list = []  # List to store content
    images = []  # List to store images
    related_keywords = []  # List to store related keywords
    
    # Prioritize Byju's URL
    response = requests.get(byjus_url)
    if response.status_code == 200:
        soup = BeautifulSoup(response.content, 'html.parser')
        content_div = soup.find('div', {'id': 'gs_res_ccl_mid'})
        if content_div:
            # Extract paragraphs
            paragraphs = content_div.find_all('p', limit=25)
            content_list.extend([p.text.strip() for p in paragraphs if p.text.strip()])
            
            # Extract images
            for img in content_div.find_all('img', limit=5):
                img_url = img.get('src')
                if img_url:
                    images.append(f"https:{img_url}")
            
            # If content found, summarize it and return
            if len(content_list) > 0:
                summarized_content = summarize_content(content_list, query_type)
                return {
                    "content": summarized_content,
                    "images": images,
                    "related_keywords": related_keywords[:3]
                }

    # If Byju's content not found or not relevant, proceed with other sources
    for url in other_sources:
        response = requests.get(url)
        if response.status_code == 200:
            soup = BeautifulSoup(response.content, 'html.parser')
            # Find content div based on source
            if "scholar.google" in url:
                content_div = soup.find('div', {'id': 'gs_res_ccl_mid'}) 
            elif "wikipedia" in url:
                content_div = soup.find('div', {'id': 'mw-content-text'})         
            elif "britannica" in url:
                content_div = soup.find('div', {'class': 'topic-identifier'})
                
            if content_div:
                # Extract paragraphs
                paragraphs = content_div.find_all('p', limit=25)
                content_list.extend([p.text.strip() for p in paragraphs if p.text.strip()])
                
                # Extract images
                for img in content_div.find_all('img', limit=5):
                    img_url = img.get('src')
                    if img_url:
                        images.append(f"https:{img_url}")
                
                # If content found, break the loop
                if len(content_list) > 0:
                    break
                else:
                    related_keywords = extract_keywords(soup.title.string)
    
    # If no content found from any source, return a message
    if not content_list:
        return {
            "content": "Content not found.",
            "images": images,
            "related_keywords": related_keywords[:3]
        }
    
    # Summarize content and return
    summarized_content = summarize_content(content_list, query_type)
    
    result = {
        "content": summarized_content,
        "images": images,
        "related_keywords": related_keywords[:3]
    }
    
    return result

# Function to extract keywords from a query
def extract_keywords(query):
    return re.findall(r'\b\w+\b', query)

# Function to summarize content based on query type
def summarize_content(content_list, query_type):
    if not content_list:
        return "Content not found."
    
    if query_type == "definition":
        return content_list[0] if content_list else "Definition not found."
    elif query_type == "essay":
        return ' '.join(content_list[:300]) if content_list else "Essay content not found."
    elif query_type == "analysis":
        return ' '.join(content_list[:25]) if content_list else "Analysis not found."
    elif query_type == "description":
        return ' '.join(content_list[:10]) if content_list else "Description not found."
    else:
        return "Invalid query type."

# Function to create PDF
def create_pdf(query, query_type, result):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    title_text = f"{query_type.capitalize()} of {query}"
    content_text = result['content']
    
    try:
        pdf.cell(200, 10, txt=title_text, ln=True, align='C')
        pdf.multi_cell(0, 10, txt=content_text)
        
        # Add images if available
        if result['images']:
            for img in result['images']:
                pdf.image(img, w=100)
        
        pdf.output("output.pdf")
        print("PDF created successfully.")
    except Exception as e:
        print(f"An error occurred while creating the PDF: {e}")

# Function to create PowerPoint presentation
def create_ppt(query, query_type, result):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"{query_type.capitalize()} of {query}"
    
    left_inch = Inches(1)
    top_inch = Inches(1.2)
    width_inch = Inches(8)
    height_inch = Inches(1.5)
    
    textbox = slide.shapes.add_textbox(left_inch, top_inch, width_inch, height_inch)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.add_paragraph()
    p.text = result['content']
    
    # Add images if available
    if result['images']:
        for img in result['images']:
            slide.shapes.add_picture(img, Inches(1), Inches(2), height=Inches(1.5))
    
    prs.save('Nobestudy.pptx')

# Function to create DOCX document
def create_docx(query, query_type, result):
    doc = Document()
    doc.add_heading(f"{query_type.capitalize()} of {query}", level=1)
    doc.add_paragraph(result['content'])
    
    # Add images if available
    if result['images']:
        for img in result['images']:
            doc.add_picture(img, width=Inches(2))
    
    doc.save('Nobestudy.docx')

# Route for home page
@app.route('/')
def home():
    return render_template('scrpe.html')

# Route for handling search request
@app.route('/search', methods=['POST'])
def search():
    query = request.form.get('query')
    query_type = request.form.get('query_type')
    
    # Check if query and query type are provided
    if not query or not query_type:
        return render_template('scrpe.html', error="Please provide both query and query type.")
    
    # Extract keywords from query
    keywords = extract_keywords(query)
    # Scrape information based on query and query type
    result = scrape_information(' '.join(keywords), query_type)
    
    return render_template('scrpe.html', query=query, query_type=query_type, result=result)

# Route for downloading files
@app.route('/download/<file_type>', methods=['POST'])
def download(file_type):
    query = request.form.get('query')
    query_type = request.form.get('query_type')
    content = request.form.get('content')
    
    # Check if query, query type, and content are provided
    if not query or not query_type or not content:
        return render_template('scrpe.html', error="Missing data to generate file.")
    
    result = {"content": content, "images": []}
    
    # Generate file based on file type
    if file_type == 'pdf':
        create_pdf(query, query_type, result)
        return send_file("output.pdf", as_attachment=True)
    elif file_type == 'ppt':
        create_ppt(query, query_type, result)
        return send_file("Nobestudy.pptx", as_attachment=True)
    elif file_type == 'docx':
        create_docx(query, query_type, result)
        return send_file("Nobestudy.docx", as_attachment=True)
    else:
        return render_template('scrpe.html', error="Invalid file type.")

if __name__ == "__main__":
    app.run()
